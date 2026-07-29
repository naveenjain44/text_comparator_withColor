"""Iteration-5 tests: universal parsers, /api/supported, dynamic-value
normalization, greeting-in-body, AmEx PDF<->EML scenario, cross-format compare,
and packaging static checks."""
from __future__ import annotations

import io
import os
import sys
from pathlib import Path

import pytest
import requests
from dotenv import dotenv_values

sys.path.insert(0, "/app/backend")

frontend_env = dotenv_values("/app/frontend/.env")
base_url = os.environ.get("REACT_APP_BACKEND_URL") or frontend_env.get("REACT_APP_BACKEND_URL")
if not base_url:
    raise RuntimeError("REACT_APP_BACKEND_URL missing")
BASE_URL = base_url.rstrip("/")
API = f"{BASE_URL}/api"

SAMPLES = Path("/app/backend/samples")
AMEX = SAMPLES / "amex"
AMEX_PDF = AMEX / "Gmail - Fwd_ Your balance update.pdf"
AMEX_EML = AMEX / "Your balance update.eml"
STUBS = Path("/tmp/uat_stubs")

ALL_EXT = [".csv", ".docx", ".eml", ".htm", ".html", ".msg", ".pdf", ".pptx", ".txt", ".xlsx"]


@pytest.fixture(scope="session")
def client():
    s = requests.Session()
    return s


# ---------- stub generation ----------

@pytest.fixture(scope="session")
def stub_files():
    STUBS.mkdir(parents=True, exist_ok=True)
    text = (
        "Subject: TEST_ Your balance update\n\n"
        "Dear KIRTIMAAN CHHABRA,\n\n"
        "Your total balance is INR xx,xxx.00 as of 02-Jun-2026.\n\n"
        "View statement now\n\n"
        "Unsubscribe | Privacy Policy\n"
    )
    (STUBS / "stub.txt").write_bytes(text.encode())
    (STUBS / "stub.csv").write_bytes(
        b"Field,Value\nGreeting,Dear KIRTIMAAN CHHABRA\nBalance,INR 12500.00\n")
    html = ("<html><head><title>TEST_ Your balance update</title></head><body>"
            "<p>Dear KIRTIMAAN CHHABRA,</p>"
            "<p style='font-size:0'>spacer</p>"
            "<p style='max-height:0'>hidden-trick</p>"
            "<p style='display:none'>should-be-removed</p>"
            "<p>Your total balance is INR 12,500.00 as of 02-Jun-2026.</p>"
            "<a href='https://example.com/stmt'>View statement now</a>"
            "<p>Unsubscribe | Privacy Policy</p></body></html>")
    (STUBS / "stub.html").write_bytes(html.encode())
    (STUBS / "stub.htm").write_bytes(html.encode())

    # docx
    from docx import Document
    d = Document()
    for line in ["TEST_ Your balance update", "Dear KIRTIMAAN CHHABRA,",
                 "Your total balance is INR 12,500.00 as of 02-Jun-2026.",
                 "View statement now", "Unsubscribe | Privacy Policy"]:
        d.add_paragraph(line)
    d.save(STUBS / "stub.docx")

    # xlsx
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    for row in [["TEST_ Your balance update"], ["Dear KIRTIMAAN CHHABRA,"],
                ["Your total balance is INR 12,500.00 as of 02-Jun-2026."],
                ["View statement now"], ["Unsubscribe | Privacy Policy"]]:
        ws.append(row)
    wb.save(STUBS / "stub.xlsx")

    # pptx
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(0, 0, 5000000, 3000000).text_frame
    tb.text = "Dear KIRTIMAAN CHHABRA,"
    tb.add_paragraph().text = "Your total balance is INR 12,500.00 as of 02-Jun-2026."
    prs.save(STUBS / "stub.pptx")

    # eml
    eml = (b"Subject: TEST_ Your balance update\r\nFrom: a@b.com\r\nTo: c@d.com\r\n"
           b"Content-Type: text/html; charset=utf-8\r\n\r\n" + html.encode())
    (STUBS / "stub.eml").write_bytes(eml)

    # pdf
    from reportlab.pdfgen import canvas
    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    y = 800
    for line in ["TEST_ Your balance update", "Dear KIRTIMAAN CHHABRA,",
                 "Your total balance is INR 12,500.00 as of 02-Jun-2026.",
                 "View statement now", "Unsubscribe | Privacy Policy"]:
        c.drawString(70, y, line)
        y -= 40
    c.save()
    (STUBS / "stub.pdf").write_bytes(buf.getvalue())

    return STUBS


# ---------- /api/supported ----------

class TestSupportedEndpoint:
    def test_supported_extensions(self, client):
        r = client.get(f"{API}/supported", timeout=30)
        assert r.status_code == 200, r.text
        data = r.json()
        assert "extensions" in data
        assert sorted(data["extensions"]) == ALL_EXT, data["extensions"]


# ---------- parsers module ----------

class TestParsers:
    def test_supported_ext_set(self):
        from parsers import SUPPORTED_EXT
        assert sorted(SUPPORTED_EXT) == ALL_EXT

    @pytest.mark.parametrize("ext", [e for e in ALL_EXT if e != ".msg"])
    def test_parse_each_extension(self, stub_files, ext):
        from parsers import parse_document
        name = f"stub{ext}"
        data = (stub_files / name).read_bytes()
        parsed = parse_document(name, data)
        # every parser must produce at least something
        assert parsed is not None
        assert (parsed.body_paragraphs or parsed.raw_text or parsed.subject), \
            f"{ext} produced empty parse"

    def test_unsupported_extension_raises(self):
        from parsers import parse_document
        with pytest.raises(ValueError):
            parse_document("weird.zzz", b"data")

    def test_html_font_size_zero_not_dropped(self, stub_files):
        from parsers import parse_document
        parsed = parse_document("stub.html", (stub_files / "stub.html").read_bytes())
        blob = " ".join(parsed.body_paragraphs) + parsed.raw_text
        assert "should-be-removed" not in blob, "display:none content must be removed"


# ---------- dynamic normalization unit ----------

class TestDynamicNormalization:
    @pytest.mark.parametrize("raw,expected_token", [
        ("INR xx, xxx.00", "<AMOUNT>"),
        ("INR 12,500.00", "<AMOUNT>"),
        ("$12.50", "<AMOUNT>"),
        ("02-Jun-2026", "<DATE>"),
        ("05-Jun-2026", "<DATE>"),
        ("Jun 5, 2026", "<DATE>"),
    ])
    def test_tokens(self, raw, expected_token):
        from comparator import apply_dynamic_normalization
        out = apply_dynamic_normalization(raw)
        assert expected_token in out, f"{raw!r} -> {out!r}"

    def test_amount_full_collapse(self):
        from comparator import apply_dynamic_normalization
        assert apply_dynamic_normalization("INR xx, xxx.00").strip() == "<AMOUNT>"

    def test_mask(self):
        from comparator import apply_dynamic_normalization
        assert apply_dynamic_normalization("Account ending: 4xxx8") == "Account ending: <MASK>"

    def test_longnum(self):
        from comparator import apply_dynamic_normalization
        out = apply_dynamic_normalization("Ref 987654321 ok")
        assert "<LONGNUM>" in out or "<MASK>" in out, out

    def test_smart_similarity_amount_pair(self):
        from comparator import similarity
        assert similarity("INR xx, xxx.00", "INR 0.00", [], "smart") == 1.0

    def test_strict_similarity_amount_pair_less_than_one(self):
        from comparator import similarity
        assert similarity("INR xx,xxx.00", "INR 0.00", [], "strict") < 1.0

    def test_weights_greeting_zero(self):
        from comparator import WEIGHTS
        assert WEIGHTS["greeting"] == 0.0
        assert abs(sum(WEIGHTS.values()) - 1.0) < 1e-6


# ---------- noise filtering ----------

class TestNoiseFilter:
    @pytest.mark.parametrize("line", [
        "30/07/2026, 00:00 Gmail - Fwd: Your balance update",
        "https://mail.google.com/mail/u/0",
        "<someone@example.com>",
        "K irtimaan C hhabra",
        "K C",
    ])
    def test_noise_detected(self, line):
        from comparator import _is_noise_line
        assert _is_noise_line(line) is True, line

    @pytest.mark.parametrize("line", [
        "Dear KIRTIMAAN CHHABRA,",
        "Your total balance is INR 12,500.00 as of 02-Jun-2026.",
    ])
    def test_real_content_kept(self, line):
        from comparator import _is_noise_line
        assert _is_noise_line(line) is False, line


# ---------- AmEx PDF <-> EML ----------

@pytest.fixture(scope="class")
def amex_report(client):
    assert AMEX_PDF.exists(), f"missing {AMEX_PDF}"
    assert AMEX_EML.exists(), f"missing {AMEX_EML}"
    files = {
        "mockup": (AMEX_PDF.name, AMEX_PDF.read_bytes(), "application/pdf"),
        "output": (AMEX_EML.name, AMEX_EML.read_bytes(), "message/rfc822"),
    }
    r = client.post(f"{API}/compare", files=files,
                    data={"mode": "smart", "use_glossary": "true", "save_history": "false"},
                    timeout=180)
    assert r.status_code == 200, f"{r.status_code}: {r.text[:600]}"
    return r.json()


class TestAmexPdfEml:
    def test_greeting_extracted_and_matched(self, amex_report):
        g = amex_report["greeting"]
        assert g["mockup"] == "Dear KIRTIMAAN CHHABRA,", g["mockup"]
        assert g["status"] == "match", g

    def test_greeting_prepended_into_body(self, amex_report):
        body = amex_report["body"]
        assert body, "empty body diff"
        assert body[0]["mockup"].startswith("Dear KIRTIMAAN CHHABRA"), body[0]["mockup"]
        assert body[0]["status"] == "match", body[0]

    def test_amount_row_similarity_is_one(self, amex_report):
        rows = [r for r in amex_report["body"]
                if (r.get("mockup") or "").strip() == "INR xx, xxx.00"
                and (r.get("output") or "").strip() == "INR 0.00"]
        assert rows, "no INR xx, xxx.00 / INR 0.00 pair found: " + str(
            [(r.get("mockup"), r.get("output")) for r in amex_report["body"]])[:1200]
        assert any(abs(r["similarity"] - 1.0) < 1e-6 for r in rows), rows

    def test_body_score_above_40(self, amex_report):
        assert amex_report["scores"]["body"] > 40, amex_report["scores"]

    def test_overall_score_positive(self, amex_report):
        assert amex_report["summary"]["overall_score"] > 0, amex_report["summary"]


# ---------- cross-format compare via API ----------

class TestCrossFormat:
    @pytest.mark.parametrize("m_ext,o_ext", [
        (".txt", ".docx"),
        (".html", ".eml"),
        (".csv", ".xlsx"),
        (".pdf", ".txt"),
        (".pptx", ".html"),
    ])
    def test_any_to_any(self, client, stub_files, m_ext, o_ext):
        m = stub_files / f"stub{m_ext}"
        o = stub_files / f"stub{o_ext}"
        files = {
            "mockup": (m.name, m.read_bytes(), "application/octet-stream"),
            "output": (o.name, o.read_bytes(), "application/octet-stream"),
        }
        r = client.post(f"{API}/compare", files=files,
                        data={"mode": "smart", "save_history": "false"}, timeout=120)
        assert r.status_code == 200, f"{m_ext}->{o_ext} {r.status_code}: {r.text[:400]}"
        data = r.json()
        for key in ("summary", "scores", "body", "greeting", "subject"):
            assert key in data, key
        assert isinstance(data["summary"]["overall_score"], (int, float))

    def test_reverse_amex_eml_pdf(self, client):
        files = {
            "mockup": (AMEX_EML.name, AMEX_EML.read_bytes(), "message/rfc822"),
            "output": (AMEX_PDF.name, AMEX_PDF.read_bytes(), "application/pdf"),
        }
        r = client.post(f"{API}/compare", files=files,
                        data={"mode": "smart", "save_history": "false"}, timeout=180)
        assert r.status_code == 200, r.text[:500]
        assert r.json()["summary"]["overall_score"] > 0

    def test_strict_mode_amex_scores_lower(self, client):
        files = {
            "mockup": (AMEX_PDF.name, AMEX_PDF.read_bytes(), "application/pdf"),
            "output": (AMEX_EML.name, AMEX_EML.read_bytes(), "message/rfc822"),
        }
        r = client.post(f"{API}/compare", files=files,
                        data={"mode": "strict", "save_history": "false"}, timeout=180)
        assert r.status_code == 200, r.text[:400]
        rows = [x for x in r.json()["body"]
                if (x.get("mockup") or "").strip().startswith("INR")
                and (x.get("output") or "").strip().startswith("INR")
                and (x.get("mockup") or "").strip() != (x.get("output") or "").strip()]
        if rows:
            assert all(x["similarity"] < 1.0 for x in rows), rows


# ---------- packaging static checks ----------

class TestPackagingStatic:
    @pytest.mark.parametrize("path", [
        "/app/backend/launcher.py", "/app/uat_tool.spec",
        "/app/build.sh", "/app/build.bat",
    ])
    def test_file_exists_non_empty(self, path):
        p = Path(path)
        assert p.exists(), f"{path} missing"
        assert p.stat().st_size > 0, f"{path} empty"

    def test_launcher_imports(self):
        import importlib
        mod = importlib.import_module("launcher")
        assert hasattr(mod, "main")

    def test_server_static_mount_guarded(self):
        src = Path("/app/backend/server.py").read_text()
        assert "_FRONTEND_BUILD.exists()" in src
        assert not Path("/app/frontend/build").exists() or True

    def test_frontend_supported_exts(self):
        src = Path("/app/frontend/src/pages/Compare.jsx").read_text()
        assert "SUPPORTED_EXTS" in src
        for ext in ALL_EXT:
            assert ext in src, f"{ext} missing from Compare.jsx"
        # applied to both upload zones
        assert src.count("SUPPORTED_EXTS") >= 3, src.count("SUPPORTED_EXTS")

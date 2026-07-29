"""
Universal file parser — extracts (subject, greeting, body, cta, footer, images, links)
from any of the supported document types:

  .docx  .eml  .msg  .htm  .html  .pdf  .txt  .pptx  .xlsx  .csv

All parsers return the same `ParsedEmail` dataclass so the comparator can be
run on any pair of files (docx↔pdf, eml↔pptx, xlsx↔txt, etc.).
"""
from __future__ import annotations

import csv as _csv
import io
import re
from pathlib import Path
from typing import List

from bs4 import BeautifulSoup

from comparator import (
    ParsedEmail,
    ImageItem,
    LinkItem,
    normalize_ws,
    parse_docx,
    parse_eml,
    parse_msg,
    _split_sections,
    _leaf_blocks,
    _parse_html_body,
)


def _wrap_text_as_email(text: str, images: List[ImageItem] | None = None,
                        links: List[LinkItem] | None = None,
                        source_subject: str = "") -> ParsedEmail:
    """Split a plain text blob into paragraphs + section-detect."""
    result = ParsedEmail()
    if source_subject:
        result.subject = normalize_ws(source_subject)
    if images:
        result.images.extend(images)
    if links:
        result.links.extend(links)

    # Extract URLs from plain text as links
    url_re = re.compile(r"https?://[^\s)>\]\"']+")
    for u in url_re.findall(text):
        if not any(l.url == u for l in result.links):
            result.links.append(LinkItem(text=u, url=u))

    # Paragraphs are separated by blank line OR single newline (whichever is present)
    paragraphs = []
    if "\n\n" in text:
        chunks = re.split(r"\n\s*\n", text)
    else:
        chunks = text.split("\n")
    for c in chunks:
        t = normalize_ws(c)
        if t:
            paragraphs.append(t)

    _split_sections(result, paragraphs)
    result.raw_text = "\n\n".join(paragraphs)
    return result


# ---------- HTML / HTM ----------

def parse_html_file(data: bytes) -> ParsedEmail:
    result = ParsedEmail()
    try:
        html = data.decode("utf-8", errors="replace")
    except Exception:
        html = data.decode("latin-1", errors="replace")
    _parse_html_body(html, result)
    # subject best-effort from <title>
    soup = BeautifulSoup(html, "lxml")
    title = soup.find("title")
    if title and title.get_text(strip=True):
        result.subject = normalize_ws(title.get_text(strip=True))
    return result


# ---------- Plain text ----------

def parse_txt(data: bytes) -> ParsedEmail:
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception:
        text = data.decode("latin-1", errors="replace")
    # First line "Subject: ..." style if present
    result = _wrap_text_as_email(text)
    return result


# ---------- CSV ----------

def parse_csv(data: bytes) -> ParsedEmail:
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception:
        text = data.decode("latin-1", errors="replace")
    rows = []
    reader = _csv.reader(io.StringIO(text))
    for row in reader:
        line = " | ".join(cell.strip() for cell in row if cell and cell.strip())
        if line:
            rows.append(line)
    return _wrap_text_as_email("\n".join(rows))


# ---------- XLSX ----------

def parse_xlsx(data: bytes) -> ParsedEmail:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    paragraphs = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return _wrap_text_as_email("\n".join(paragraphs))


# ---------- PPTX ----------

def parse_pptx(data: bytes) -> ParsedEmail:
    from pptx import Presentation
    from pptx.util import Emu  # noqa
    prs = Presentation(io.BytesIO(data))
    lines: List[str] = []
    images: List[ImageItem] = []
    links: List[LinkItem] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    parts = []
                    for run in para.runs:
                        parts.append(run.text or "")
                        try:
                            hl = run.hyperlink.address
                            if hl:
                                links.append(LinkItem(text=(run.text or "").strip(), url=hl))
                        except Exception:
                            pass
                    txt = normalize_ws("".join(parts))
                    if txt:
                        lines.append(txt)
            # picture
            if getattr(shape, "shape_type", None) == 13:  # PICTURE
                name = getattr(shape, "name", None) or "image"
                alt = ""
                try:
                    # pptx puts descr on the shape's XML nvPicPr/cNvPr @descr
                    descr = shape._element.xpath(".//p:nvPicPr/p:cNvPr/@descr")
                    if descr:
                        alt = descr[0]
                except Exception:
                    pass
                images.append(ImageItem(filename=name, alt=alt))
    return _wrap_text_as_email("\n".join(lines), images=images, links=links)


# ---------- PDF ----------

def parse_pdf(data: bytes) -> ParsedEmail:
    import pdfplumber
    paragraphs: List[str] = []
    links: List[LinkItem] = []
    subject_guess = ""
    _subj_re = re.compile(r"^\s*(?:fwd:|re:|subject:)\s*(.+)$", re.I)
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            try:
                lines = page.extract_text_lines(strip=True) or []
            except Exception:
                lines = []
            if lines:
                current: List[str] = []
                prev_bottom = None
                heights = [l.get("bottom", 0) - l.get("top", 0) for l in lines if l.get("bottom") and l.get("top")]
                avg_h = (sum(heights) / len(heights)) if heights else 12
                for l in lines:
                    text = normalize_ws(l.get("text") or "")
                    if not text:
                        continue
                    if not subject_guess:
                        m = _subj_re.match(text)
                        if m:
                            subject_guess = m.group(1).strip()
                    top = l.get("top", 0)
                    if prev_bottom is not None and (top - prev_bottom) > avg_h * 0.9:
                        if current:
                            paragraphs.append(normalize_ws(" ".join(current)))
                            current = []
                    current.append(text)
                    prev_bottom = l.get("bottom", top)
                if current:
                    paragraphs.append(normalize_ws(" ".join(current)))
            else:
                text = page.extract_text() or ""
                for block in re.split(r"\n\s*\n+", text):
                    joined = " ".join(x.strip() for x in block.split("\n") if x.strip())
                    joined = normalize_ws(joined)
                    if joined:
                        if not subject_guess:
                            m = _subj_re.match(joined)
                            if m:
                                subject_guess = m.group(1).strip()
                        paragraphs.append(joined)
            try:
                for a in (page.hyperlinks or []):
                    uri = a.get("uri")
                    if uri:
                        links.append(LinkItem(text=uri, url=uri))
            except Exception:
                pass
    text_all = "\n\n".join(paragraphs)
    return _wrap_text_as_email(text_all, links=links, source_subject=subject_guess)


# ---------- Dispatcher ----------

SUPPORTED_EXT = {
    ".docx", ".eml", ".msg", ".htm", ".html", ".pdf",
    ".txt", ".pptx", ".xlsx", ".csv",
}


def parse_document(filename: str, data: bytes) -> ParsedEmail:
    ext = Path(filename or "").suffix.lower()
    if ext == ".docx":
        return parse_docx(data)
    if ext == ".eml":
        return parse_eml(data)
    if ext == ".msg":
        return parse_msg(data)
    if ext in (".htm", ".html"):
        return parse_html_file(data)
    if ext == ".pdf":
        return parse_pdf(data)
    if ext == ".txt":
        return parse_txt(data)
    if ext == ".pptx":
        return parse_pptx(data)
    if ext == ".xlsx":
        return parse_xlsx(data)
    if ext == ".csv":
        return parse_csv(data)
    raise ValueError(
        f"Unsupported file type '{ext}'. Supported: {sorted(SUPPORTED_EXT)}"
    )
